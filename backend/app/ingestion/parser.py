from pathlib import Path

from app.ingestion.clean import clean_chunk


SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".html", ".pptx", ".docx", ".png", ".jpg", ".jpeg", ".tiff", ".webp"}


def parse_document(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext in {".png", ".jpg", ".jpeg", ".tiff", ".webp"}:
        return _parse_image(file_path)
    elif ext == ".txt":
        return _parse_txt(file_path)
    elif ext == ".html":
        return _parse_html(file_path)
    elif ext == ".pptx":
        return _parse_pptx(file_path)
    elif ext == ".docx":
        return _parse_docx(file_path)
    raise ValueError(f"Unsupported file type: {ext}")


MIN_PAGE_WORDS = 20
SCANNED_COVERAGE = 0.6
MIN_IMAGE_CROP = 50 * 50  # skip tiny icons/logos (area in pt^2)


def _parse_pdf(file_path: str) -> str:
    """Per-page PDF parsing.

    Decision per page using the text layer word count and embedded-image coverage:
      - text layer present, no images        -> keep text (fast path)
      - text layer present, images embedded  -> keep text + OCR image crops
      - no text, image fills the page (scan) -> OCR the whole page render
      - no text, no image                    -> keep thin text (cleaner may drop it)
    """
    import pymupdf
    from app.core.ocr import ocr_image

    doc = pymupdf.open(file_path)
    parts = []
    for i, page in enumerate(doc, start=1):
        page_text = _parse_page(page, ocr_image)
        if page_text:
            parts.append(f"--- Page {i} ---\n{page_text}")
    return "\n\n".join(parts)


def _parse_page(page, ocr_image) -> str:
    text = page.get_text() or ""
    words = len(text.split())
    images = page.get_image_info()

    page_area = page.rect.width * page.rect.height or 1
    coverage = 0.0
    crops = []
    for img in images:
        x0, y0, x1, y1 = img["bbox"]
        area = (x1 - x0) * (y1 - y0)
        coverage += area
        if area >= MIN_IMAGE_CROP:
            crops.append((x0, y0, x1, y1))

    if words >= MIN_PAGE_WORDS:
        # Text layer present. OCR embedded images if any.
        if crops:
            for x0, y0, x1, y1 in crops:
                pix = page.get_pixmap(clip=pymupdf.Rect(x0, y0, x1, y1), dpi=200)
                img_text = ocr_image(pix.tobytes("png"))
                if img_text.strip():
                    text += f"\n[Image: {img_text}]"
        return text.strip()

    # Thin/no text layer
    if coverage / page_area > SCANNED_COVERAGE:
        # Image fills the page -> scanned page, OCR the whole render
        pix = page.get_pixmap(dpi=200)
        ocr_text = ocr_image(pix.tobytes("png"))
        return ocr_text.strip() or text.strip()
    return text.strip()


def _parse_image(file_path: str) -> str:
    from app.core.ocr import ocr_image
    with open(file_path, "rb") as f:
        return ocr_image(f.read())


def _parse_txt(file_path: str) -> str:
    with open(file_path) as f:
        return f.read()


def _parse_html(file_path: str) -> str:
    from bs4 import BeautifulSoup
    with open(file_path) as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    return soup.get_text(separator="\n")


def _parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        slide_text = "\n".join(t for t in texts if t)
        if slide_text.strip():
            parts.append(f"--- Slide {i} ---\n{slide_text.strip()}")
    return "\n\n".join(parts)


def _parse_docx(file_path: str) -> str:
    import re
    from docx import Document
    doc = Document(file_path)
    lines = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        style = p.style.name or ""
        match = re.match(r"Heading\s*(\d+)", style)
        if match:
            text = "#" * int(match.group(1)) + " " + text
        lines.append(text)
    return "\n".join(lines)


def parse_document_clean(file_path: str) -> str:
    """Parse then run the deterministic cleaning filters, returning clean text (or "" if dropped)."""
    text = parse_document(file_path)
    return clean_chunk(text)
